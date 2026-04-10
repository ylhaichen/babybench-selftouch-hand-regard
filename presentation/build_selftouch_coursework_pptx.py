from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
ASSETS_DIR = PRESENTATION_DIR / "assets"
OUTPUT_PPTX = PRESENTATION_DIR / "COMP0223_selftouch_video_deck.pptx"


def add_full_slide_image(slide, image_path: Path, slide_width, slide_height) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=slide_width, height=slide_height)


def add_textbox(slide, left, top, width, height, text, font_size, bold=False, color=(24, 48, 84), align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    return box


def add_bullets(slide, left, top, width, lines):
    box = slide.shapes.add_textbox(left, top, width, Inches(2.8))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.bullet = True
        p.space_after = Pt(12)
        for run in p.runs:
            run.font.name = "Aptos"
            run.font.size = Pt(22)
            run.font.color.rgb = RGBColor(43, 62, 84)


def build_slide_1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(250, 252, 255)

    add_textbox(
        slide,
        Inches(0.65),
        Inches(0.55),
        Inches(5.8),
        Inches(1.2),
        "Learning Infant-Like Self-Touch in BabyBench",
        28,
        bold=True,
        color=(15, 39, 68),
    )
    add_textbox(
        slide,
        Inches(0.67),
        Inches(1.45),
        Inches(5.0),
        Inches(0.6),
        "COMP0223 Robot Learning Coursework",
        15,
        color=(91, 109, 132),
    )
    add_bullets(
        slide,
        Inches(0.7),
        Inches(2.0),
        Inches(4.8),
        [
            "Task: self-touch learning without extrinsic rewards",
            "Platform: BabyBench + MIMo",
            "Approach: intrinsic-motivation RL",
        ],
    )

    hero = slide.shapes.add_picture(str(ASSETS_DIR / "slide1_hero_frame.png"), Inches(6.2), Inches(1.2), width=Inches(6.2))
    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        hero.left - Inches(0.08),
        hero.top - Inches(0.08),
        hero.width + Inches(0.16),
        hero.height + Inches(0.16),
    )
    border.fill.background()
    border.line.color.rgb = RGBColor(184, 204, 228)
    border.line.width = Pt(2.5)
    slide.shapes._spTree.remove(border._element)
    slide.shapes._spTree.insert(2, border._element)

    add_textbox(
        slide,
        Inches(0.7),
        Inches(5.85),
        Inches(11.6),
        Inches(0.65),
        "Question: can the robot discover its own body through self-generated exploration?",
        18,
        color=(55, 79, 109),
    )


def build_figure_slide(prs: Presentation, image_name: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_full_slide_image(slide, ASSETS_DIR / image_name, prs.slide_width, prs.slide_height)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_figure_slide(prs, "slide2_concept_graphic.png")
    build_figure_slide(prs, "slide3_pipeline_graphic.png")
    build_figure_slide(prs, "slide4_training_progress.png")
    build_figure_slide(prs, "slide5_final_results.png")

    prs.save(str(OUTPUT_PPTX))
    print(f"Saved PowerPoint deck to: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
