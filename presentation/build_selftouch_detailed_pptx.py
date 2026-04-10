from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[1]
PRESENTATION_DIR = ROOT / "presentation"
ASSETS_DIR = PRESENTATION_DIR / "assets_detailed"
OUTPUT_PPTX = PRESENTATION_DIR / "COMP0223_selftouch_video_deck_detailed.pptx"


SLIDES = [
    "slide1_cover_detailed.png",
    "slide2_motivation_detailed.png",
    "slide3_robotics_motivation.png",
    "slide4_method_detailed.png",
    "slide5_algorithm_detailed.png",
    "slide6_training_detailed.png",
    "slide7_results_detailed.png",
]


def add_full_slide_image(slide, image_path: Path, slide_width, slide_height) -> None:
    slide.shapes.add_picture(str(image_path), 0, 0, width=slide_width, height=slide_height)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for image_name in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_full_slide_image(slide, ASSETS_DIR / image_name, prs.slide_width, prs.slide_height)

    prs.save(str(OUTPUT_PPTX))
    print(f"Saved detailed PowerPoint deck to: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
